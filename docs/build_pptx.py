"""生成 Local Window Copilot 项目介绍 PPTX。

运行：
    python docs/build_pptx.py
输出：
    docs/Local_Window_Copilot.pptx
"""
from __future__ import annotations

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Emu
from pptx.oxml.ns import qn
from lxml import etree

# ---------- 配色（对齐 HTML 版） ----------
BG = RGBColor(0x0D, 0x11, 0x17)
PANEL = RGBColor(0x1C, 0x22, 0x30)
LINE = RGBColor(0x2A, 0x31, 0x42)
TEXT = RGBColor(0xE6, 0xED, 0xF3)
MUTED = RGBColor(0x9A, 0xA7, 0xB8)
ACCENT = RGBColor(0x4D, 0xD0, 0xE1)
ACCENT2 = RGBColor(0x7C, 0x9E, 0xFF)
WARN = RGBColor(0xFF, 0xB8, 0x6B)
GOOD = RGBColor(0x7E, 0xE7, 0x87)
BAD = RGBColor(0xFF, 0x7B, 0x72)
CODE_BG = RGBColor(0x0E, 0x11, 0x16)

FONT = "Microsoft YaHei"
MONO = "Consolas"

prs = Presentation()
prs.slide_width = Inches(13.333)   # 16:9
prs.slide_height = Inches(7.5)

# 设计逻辑像素 1280x720，转成 EMU
_EMU_PER_PX = prs.slide_width // 1280
def PX(v):
    """逻辑像素 → EMU"""
    return int(v * _EMU_PER_PX)

SW, SH = 1280, 720  # 用逻辑像素设计

BLANK = prs.slide_layouts[6]  # 空白版式


# ---------- 工具函数 ----------
def add_bg(slide, color=BG):
    """全屏背景矩形。"""
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, PX(SW), PX(SH))
    rect.line.fill.background()
    rect.fill.solid()
    rect.fill.fore_color.rgb = color
    rect.shadow.inherit = False
    # 置底
    spTree = rect._element.getparent()
    spTree.remove(rect._element)
    spTree.insert(2, rect._element)
    return rect


def add_text(slide, left, top, width, height, text, *,
             size=14, color=TEXT, bold=False, font=FONT,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    """纯文本框。"""
    tb = slide.shapes.add_textbox(PX(left), PX(top), PX(width), PX(height))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tb


def add_rect(slide, left, top, width, height, *,
             fill=PANEL, line=LINE, line_w=0.75):
    """矩形面板。"""
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, PX(left), PX(top), PX(width), PX(height))
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(line_w)
    sh.shadow.inherit = False
    return sh


def add_round(slide, left, top, width, height, *,
              fill=PANEL, line=LINE, line_w=0.75, radius=0.08):
    """圆角矩形卡片。"""
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, PX(left), PX(top), PX(width), PX(height))
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(line_w)
    sh.shadow.inherit = False
    try:
        sh.adjustments[0] = radius
    except Exception:
        pass
    return sh


def add_multi(slide, left, top, width, height, lines):
    """多段富文本。lines = [(text, {size, color, bold, font}), ...]
    每个元素是一段（同一段内不分 run）。"""
    tb = slide.shapes.add_textbox(PX(left), PX(top), PX(width), PX(height))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    first = True
    for text, opts in lines:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.alignment = opts.get("align", PP_ALIGN.LEFT)
        p.space_after = Pt(opts.get("space_after", 0))
        p.line_spacing = opts.get("line_spacing", 1.15)
        # 支持一个段落里多 run：
        #   形式 1：("plain text", {opts})
        #   形式 2：([(t, opts), ...], {paragraph opts})
        runs = text if isinstance(text, list) else opts.get("runs")
        if runs:
            for t, ro in runs:
                r = p.add_run()
                r.text = t
                r.font.name = ro.get("font", FONT)
                r.font.size = Pt(ro.get("size", 14))
                r.font.bold = ro.get("bold", False)
                r.font.color.rgb = ro.get("color", TEXT)
        else:
            r = p.add_run()
            r.text = text
            r.font.name = opts.get("font", FONT)
            r.font.size = Pt(opts.get("size", 14))
            r.font.bold = opts.get("bold", False)
            r.font.color.rgb = opts.get("color", TEXT)
    return tb


def Pt(p):
    """兼容 python-pptx 的 Pt（这里直接复用）"""
    from pptx.util import Pt as _Pt
    return _Pt(p)


def add_header(slide, tag, page_no, total=13):
    """页眉：左侧 tag，右侧页码。"""
    add_text(slide, 60, 32, 600, 24, tag, size=11, color=ACCENT, bold=True)
    add_text(slide, SW - 160, 32, 100, 24, f"{page_no:02d} / {total:02d}",
             size=11, color=MUTED, align=PP_ALIGN.RIGHT)


def add_title(slide, text, top=70):
    add_text(slide, 60, top, SW - 120, 50, text, size=28, color=TEXT, bold=True)


def add_lead(slide, text, top=130):
    add_text(slide, 60, top, SW - 120, 60, text, size=14, color=MUTED)


def add_code_block(slide, left, top, width, height, code, size=11):
    """代码/示意图块。"""
    add_rect(slide, left, top, width, height, fill=CODE_BG, line=LINE)
    tb = slide.shapes.add_textbox(left + 16, top + 12, width - 32, height - 24)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    lines = code.split("\n")
    first = True
    for ln in lines:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.line_spacing = 1.2
        r = p.add_run()
        r.text = ln if ln else " "
        r.font.name = MONO
        r.font.size = Pt(size)
        r.font.color.rgb = RGBColor(0xC9, 0xD4, 0xE3)


# =========================================================================
# 第 1 页 封面
# =========================================================================
def slide_cover():
    s = prs.slides.add_slide(BLANK)
    add_bg(s)
    # 网格背景（用浅色线模拟）
    for x in range(0, SW, 60):
        ln = s.shapes.add_connector(1, PX(x), 0, PX(x), PX(SH))
        ln.line.color.rgb = LINE
        ln.line.width = Pt(0.5)
    for y in range(0, SH, 60):
        ln = s.shapes.add_connector(1, 0, PX(y), PX(SW), PX(y))
        ln.line.color.rgb = LINE
        ln.line.width = Pt(0.5)

    # eyebrow
    add_text(s, 80, 200, 700, 22, "LOCAL · WINDOWS · PRIVACY-FIRST",
             size=12, color=ACCENT, bold=True)
    # 标题
    add_multi(s, 80, 240, 900, 180, [
        ("Local Window", {"size": 54, "bold": True, "color": TEXT, "line_spacing": 1.1}),
        ("Copilot", {"size": 54, "bold": True, "color": ACCENT, "line_spacing": 1.1}),
    ])
    # 副标题
    add_text(s, 80, 400, 760, 60,
             "运行在 Windows 本地桌面的 AI 伙伴：观察前台窗口、按需检索上下文、用可追溯证据回答。",
             size=16, color=TEXT)
    # meta
    add_multi(s, 80, 490, 900, 80, [
        ([("汇报人　", {"size": 12, "color": MUTED}),
          ("宋林蔚", {"size": 12, "color": TEXT, "bold": True}),
          ("　　组会　", {"size": 12, "color": MUTED}),
          ("2026-07-07", {"size": 12, "color": TEXT, "bold": True}),
          ("　　仓库　", {"size": 12, "color": MUTED}),
          ("github.com/fjnuslw", {"size": 12, "color": TEXT, "bold": True})],
         {"size": 12}),
    ])
    # 右侧吉祥物圆
    mascot = add_shape_circle(s, SW - 280, 220, 200)
    try:
        s.shapes.add_picture("assets/mascot/composed/mascot_idle.png",
                             PX(SW - 260), PX(240), width=PX(160))
    except Exception:
        add_text(s, SW - 260, 290, 160, 60, " mascot ", size=14,
                 color=ACCENT, align=PP_ALIGN.CENTER)


def add_shape_circle(s, left, top, size):
    sh = s.shapes.add_shape(MSO_SHAPE.OVAL, PX(left), PX(top), PX(size), PX(size))
    sh.fill.solid()
    sh.fill.fore_color.rgb = RGBColor(0x1B, 0x23, 0x30)
    sh.line.color.rgb = LINE
    sh.line.width = Pt(1)
    sh.shadow.inherit = False
    return sh


# =========================================================================
# 第 2 页 项目定位
# =========================================================================
def slide_position():
    s = prs.slides.add_slide(BLANK)
    add_bg(s)
    add_header(s, "01 · 项目定位", 2)
    add_title(s, "这是什么？不是什么？")
    add_lead(s, "一句话定位：运行在本地 Windows 桌面的 AI 伙伴，后台把目标窗口截图转成结构化观察，前台在用户提问时按需取证、用证据回答。")

    # 左右对比卡片
    card_w, card_h = 560, 240
    top = 220
    # 要做（绿色）
    add_round(s, 60, top, card_w, card_h, fill=RGBColor(0x12, 0x1A, 0x14),
              line=RGBColor(0x2E, 0x5A, 0x36))
    add_text(s, 84, top + 18, 200, 22, "✓ 项目要做", size=12, color=GOOD, bold=True)
    add_multi(s, 84, top + 52, card_w - 48, card_h - 60, [
        ("快速观察当前前台窗口", {"size": 13, "color": MUTED, "line_spacing": 1.5, "space_after": 4}),
        ("用本地 VLM 生成结构化观察与关键点", {"size": 13, "color": MUTED, "line_spacing": 1.5, "space_after": 4}),
        ("用户追问时结合观察与短期记忆回答", {"size": 13, "color": MUTED, "line_spacing": 1.5, "space_after": 4}),
        ("本地 RuntimeStore 保存状态与短期记忆", {"size": 13, "color": MUTED, "line_spacing": 1.5, "space_after": 4}),
        ("低打扰、可观测、可复盘、可清空", {"size": 13, "color": MUTED, "line_spacing": 1.5}),
    ])

    # 不做（红色）
    add_round(s, SW - 60 - card_w, top, card_w, card_h, fill=RGBColor(0x1A, 0x12, 0x14),
              line=RGBColor(0x5A, 0x2E, 0x2E))
    add_text(s, SW - 60 - card_w + 24, top + 18, 200, 22, "✕ 明确不做",
             size=12, color=BAD, bold=True)
    add_multi(s, SW - 60 - card_w + 24, top + 52, card_w - 48, card_h - 60, [
        ("不自动点击、输入、提交、操作电脑", {"size": 13, "color": MUTED, "line_spacing": 1.5, "space_after": 4}),
        ("不做复杂规划型自主 Agent", {"size": 13, "color": MUTED, "line_spacing": 1.5, "space_after": 4}),
        ("不做 OCR / UI Automation 主链路", {"size": 13, "color": MUTED, "line_spacing": 1.5, "space_after": 4}),
        ("不做长期人格记忆", {"size": 13, "color": MUTED, "line_spacing": 1.5, "space_after": 4}),
        ("不堆叠 fallback 掩盖主链路问题", {"size": 13, "color": MUTED, "line_spacing": 1.5}),
    ])

    # 底部产品中心卡
    add_round(s, 60, 490, SW - 120, 130, fill=PANEL, line=LINE)
    add_text(s, 84, 508, 300, 24, "产品中心", size=14, color=ACCENT, bold=True)
    add_multi(s, 84, 540, SW - 168, 70, [
        ([("陪伴式桌面伙伴；截图 + VLM 是", {"size": 13, "color": MUTED}),
          ("后台感知能力", {"size": 13, "color": TEXT, "bold": True}),
          ("，不是前台产品本体。默认不汇报屏幕，默认陪伴；只有用户邀请时才进入分析。",
           {"size": 13, "color": MUTED})],
         {"line_spacing": 1.5}),
    ])


# =========================================================================
# 第 3 页 三条主线
# =========================================================================
def slide_three_lines():
    s = prs.slides.add_slide(BLANK)
    add_bg(s)
    add_header(s, "02 · 核心思想", 3)
    add_title(s, "三条主线，职责分离")
    add_lead(s, "整个项目就是三条互不污染的链路：观察线只产生证据，对话线只承接用户问题，调试线只让一切可审计。")

    top = 220
    gap = 20
    card_w = (SW - 120 - gap * 2) // 3
    card_h = 230

    cards = [
        ("1", "观察线 · 产生证据", ACCENT,
         "目标窗口截图 → MiniCPM-V 结构化观察 → SQLite 记录 + 截图文件。只写证据，不负责对话、不主动建议、不污染长期记忆。"),
        ("2", "对话线 · 按需取证", ACCENT2,
         "用户问题 + profile + 会话历史 → 模型自己决定是否调 memory.search(query) → 拿到证据后基于证据回答。"),
        ("3", "调试线 · 全程可审计", WARN,
         "WebUI 展示 latest / history / 原始 JSON / 截图 / 工具 trace / runtime logs，回答\"模型查了什么、拿到哪条记录\"。"),
    ]
    for i, (num, t, color, body) in enumerate(cards):
        left = 60 + i * (card_w + gap)
        add_round(s, left, top, card_w, card_h, fill=PANEL, line=LINE)
        # 编号块
        add_round(s, left + 24, top + 24, 36, 30, fill=RGBColor(0x14, 0x2A, 0x2E),
                  line=None, radius=0.25)
        add_text(s, left + 24, top + 28, 36, 22, num, size=14, color=ACCENT,
                 bold=True, align=PP_ALIGN.CENTER)
        add_text(s, left + 24, top + 68, card_w - 48, 30, t, size=15, color=color, bold=True)
        add_text(s, left + 24, top + 108, card_w - 48, 110, body, size=12, color=MUTED)

    # 底部设计哲学卡
    add_round(s, 60, 480, SW - 120, 140, fill=RGBColor(0x12, 0x1F, 0x24),
              line=RGBColor(0x2E, 0x4A, 0x52))
    add_text(s, 84, 498, 300, 24, "设计哲学", size=14, color=ACCENT, bold=True)
    add_multi(s, 84, 530, SW - 168, 80, [
        ([("主链路不行就", {"size": 13, "color": MUTED}),
          ("明确失败、暂停或提示用户", {"size": 13, "color": WARN, "bold": True}),
          ("，不引入掩盖问题的替代链。缺截图、缺观察、工具失败都要", {"size": 13, "color": MUTED}),
          ("明确暴露", {"size": 13, "color": TEXT, "bold": True}),
          ("，而不是用泛泛文本兜底。", {"size": 13, "color": MUTED})],
         {"line_spacing": 1.5}),
    ])


# =========================================================================
# 第 4 页 系统架构
# =========================================================================
def slide_arch():
    s = prs.slides.add_slide(BLANK)
    add_bg(s)
    add_header(s, "03 · 系统架构", 4)
    add_title(s, "三层架构：桌面 → 后端 → 本地模型")

    diagram = (
"┌──────────────────────────────────────────────────────────────────┐\n"
"│  Desktop Floating Window                                          │\n"
"│  Idle / Observing / Analyzing / Privacy / Error  ·  5 种状态     │\n"
"└──────────────────────────────────────────────────────────────────┘\n"
"                               │  FastAPI (127.0.0.1:18081)\n"
"                               ▼\n"
"┌──────────────────────────────────────────────────────────────────┐\n"
"│  Backend   assistant.py  ·  window.py  ·  webui.py             │\n"
"│                                                                  │\n"
"│  ChatAgent      profile 冻结 → probe → stream(+tools)         │\n"
"│                 └─ memory.search(query)  ·  runtime trace        │\n"
"│                                                                  │\n"
"│  Observation pipeline                                            │\n"
"│                 ├─ window capture                                │\n"
"│                 ├─ MiniCPM-V structured observation              │\n"
"│                 └─ window summary store                          │\n"
"│                                                                  │\n"
"│  RuntimeStore   runtime_json · runtime_events · chat_history_fts │\n"
"└──────────────────────────────────────────────────────────────────┘\n"
"                               │  OpenAI-compatible API\n"
"                               ▼\n"
"┌──────────────────────────────────────────────────────────────────┐\n"
"│  llama.cpp server  +  MiniCPM-V 4.6  (256K context)         │\n"
"└──────────────────────────────────────────────────────────────────┘"
    )
    add_code_block(s, 60, 130, SW - 120, 360, diagram, size=10)

    # 底部两卡
    card_w, card_h = 600, 110
    top = 510
    add_round(s, 60, top, card_w, card_h, fill=PANEL, line=LINE)
    add_text(s, 84, top + 18, 300, 22, "关键设计", size=13, color=ACCENT, bold=True)
    add_text(s, 84, top + 48, card_w - 48, 56,
             "所有运行时 JSON 状态走 SQLite RuntimeStore；截图存文件系统，SQLite 只存路径/hash。零外部 DB 依赖。",
             size=11, color=MUTED)

    add_round(s, SW - 60 - card_w, top, card_w, card_h, fill=PANEL, line=LINE)
    add_text(s, SW - 60 - card_w + 24, top + 18, 300, 22, "进程边界",
             size=13, color=ACCENT, bold=True)
    add_text(s, SW - 60 - card_w + 24, top + 48, card_w - 48, 56,
             "桌面浮窗 ←FastAPI→ Python 后端 ←OpenAI API→ llama.cpp。三层都可独立重启，互不耦合。",
             size=11, color=MUTED)


# =========================================================================
# 第 5 页 观察线
# =========================================================================
def slide_observe():
    s = prs.slides.add_slide(BLANK)
    add_bg(s)
    add_header(s, "04 · 观察线", 5)
    add_title(s, "观察线：从截图到结构化证据")
    add_lead(s, "观察线只产生证据，不做对话。输入是窗口截图 + metadata，输出是写入 SQLite 的完整结构化字段 + 截图 PNG。")

    # 流程横向 5 步
    top = 220
    steps = [
        ("目标窗口选择", "排除 Copilot 自身窗 / 优先前台 / 兜底枚举"),
        ("截图捕获", "Win32 + 标题/进程/bounds/hash/捕获时间"),
        ("VLM 分析", "MiniCPM-V 按 prompt 契约输出结构化 JSON"),
        ("SQLite 落盘", "latest_analysis + summaries 滚动历史"),
        ("截图文件", "backend/data/captures/*.png"),
    ]
    n = len(steps)
    arrow_w = 18
    step_w = (SW - 120 - arrow_w * (n - 1)) // n
    step_h = 100
    for i, (t, body) in enumerate(steps):
        left = 60 + i * (step_w + arrow_w)
        fill = RGBColor(0x12, 0x1F, 0x24) if i == 0 else PANEL
        line = RGBColor(0x2E, 0x4A, 0x52) if i == 0 else LINE
        add_round(s, left, top, step_w, step_h, fill=fill, line=line)
        add_text(s, left + 14, top + 14, step_w - 28, 24, t, size=12, color=TEXT, bold=True)
        add_text(s, left + 14, top + 44, step_w - 28, 50, body, size=10, color=MUTED)
        if i < n - 1:
            add_text(s, left + step_w, top + 35, arrow_w, 30, "→",
                     size=20, color=ACCENT, align=PP_ALIGN.CENTER)

    # 下方两卡
    card_w, card_h = 600, 230
    top2 = 360
    add_round(s, 60, top2, card_w, card_h, fill=PANEL, line=LINE)
    add_text(s, 84, top2 + 18, 400, 24, "结构化字段（每条记录必含）", size=14, color=ACCENT, bold=True)
    add_multi(s, 84, top2 + 52, card_w - 48, card_h - 60, [
        ("summary / key_points　索引与检索 rank_text", {"size": 12, "color": MUTED, "line_spacing": 1.55, "space_after": 4}),
        ("regions / visible_text / ui_elements　区域级证据", {"size": 12, "color": MUTED, "line_spacing": 1.55, "space_after": 4}),
        ("entities / uncertain_areas　实体 + 诚实标注不确定", {"size": 12, "color": MUTED, "line_spacing": 1.55, "space_after": 4}),
        ("screenshot_path / hash / bounds　可追溯", {"size": 12, "color": MUTED, "line_spacing": 1.55, "space_after": 4}),
        ("vision_input　模型输入参数", {"size": 12, "color": MUTED, "line_spacing": 1.55}),
    ])

    add_round(s, SW - 60 - card_w, top2, card_w, card_h, fill=PANEL, line=LINE)
    add_text(s, SW - 60 - card_w + 24, top2 + 18, 400, 24, "失败分型（不掩盖问题）",
             size=14, color=ACCENT, bold=True)
    add_multi(s, SW - 60 - card_w + 24, top2 + 52, card_w - 48, card_h - 60, [
        ("capture_failed　没拿到可捕获窗口", {"size": 12, "color": MUTED, "line_spacing": 1.5, "space_after": 3}),
        ("vision_failed　VLM 请求失败", {"size": 12, "color": MUTED, "line_spacing": 1.5, "space_after": 3}),
        ("vision_length　命中输出上限且无法解析", {"size": 12, "color": MUTED, "line_spacing": 1.5, "space_after": 3}),
        ("vision_parse_failed　响应非结构化 JSON", {"size": 12, "color": MUTED, "line_spacing": 1.5, "space_after": 3}),
        ("store_failed　SQLite 写入失败", {"size": 12, "color": MUTED, "line_spacing": 1.5, "space_after": 8}),
        ([("UI 与日志必须显示真实失败点，", {"size": 11, "color": MUTED}),
          ("不允许", {"size": 11, "color": BAD, "bold": True}),
          ("用泛化回答遮盖。", {"size": 11, "color": MUTED})],
         {"line_spacing": 1.5}),
    ])


# =========================================================================
# 第 6 页 对话线 probe->stream
# =========================================================================
def slide_dialog():
    s = prs.slides.add_slide(BLANK)
    add_bg(s)
    add_header(s, "05 · 对话线", 6)
    add_title(s, "对话线：probe → stream 两段式")
    add_lead(s, "小 VLM 的 function calling 不稳定——有时会把\"调用 memory.search\"写成自然语言而非结构化 tool_calls。两段式把这个不稳定性限制在工具决策层。")

    diagram = (
"1. probe 阶段   (complete_chat_response + tools)\n"
"   模型决定是否调 memory.search(query)\n"
"   → 只读取结构化 tool_calls，content 一律丢弃\n"
"\n"
"2. 工具执行     后端执行检索，生成带 record_id/source 的结果，追加到 messages\n"
"\n"
"3. stream 阶段 (stream_chat + tools)\n"
"   流式生成最终答案；tools 始终可用，模型仍可在流中发起 tool_calls\n"
"\n"
"4. 流式工具循环  stream 中检测到 tool_calls:\n"
"   暂停 yield → 执行工具 → 追加结果 → 重新 stream → 继续产出答案\n"
"   最多 2 轮，防止无限循环"
    )
    add_code_block(s, 60, 200, SW - 120, 240, diagram, size=11)

    # 下方对比
    card_w, card_h = 600, 130
    top = 470
    add_round(s, 60, top, card_w, card_h, fill=RGBColor(0x12, 0x1A, 0x14),
              line=RGBColor(0x2E, 0x5A, 0x36))
    add_text(s, 84, top + 16, 300, 22, "✓ 这样做解决了", size=12, color=GOOD, bold=True)
    add_multi(s, 84, top + 48, card_w - 48, card_h - 56, [
        ("probe 失败时 stream 仍有第二次机会调工具", {"size": 12, "color": MUTED, "line_spacing": 1.45, "space_after": 3}),
        ("不会把\"我要调用 memory.search\"这句自然语言当答案返回", {"size": 12, "color": MUTED, "line_spacing": 1.45, "space_after": 3}),
        ("职责清晰：probe 决策，stream 产出", {"size": 12, "color": MUTED, "line_spacing": 1.45}),
    ])

    add_round(s, SW - 60 - card_w, top, card_w, card_h, fill=RGBColor(0x1A, 0x12, 0x14),
              line=RGBColor(0x5A, 0x2E, 0x2E))
    add_text(s, SW - 60 - card_w + 24, top + 16, 300, 22, "✕ 禁止回归",
             size=12, color=BAD, bold=True)
    add_multi(s, SW - 60 - card_w + 24, top + 48, card_w - 48, card_h - 56, [
        ("把 probe 的 content 当最终答案直接返回", {"size": 12, "color": MUTED, "line_spacing": 1.45, "space_after": 3}),
        ("独立 planner / orchestrator 模型先规划再回答", {"size": 12, "color": MUTED, "line_spacing": 1.45, "space_after": 3}),
        ("默认把 window:latest_analysis 塞进 prompt", {"size": 12, "color": MUTED, "line_spacing": 1.45}),
    ])


# =========================================================================
# 第 7 页 唯一工具
# =========================================================================
def slide_tool():
    s = prs.slides.add_slide(BLANK)
    add_bg(s)
    add_header(s, "06 · 工具设计", 7)
    add_title(s, "唯一工具：memory.search(query)")
    add_lead(s, "模型可见工具只有一个。它是证据检索工具，不是固定上下文拼包。返回内容必须与 query 相关，并标明来源。")

    # 左卡：可检索来源
    card_w, card_h = 600, 380
    top = 200
    add_round(s, 60, top, card_w, card_h, fill=PANEL, line=LINE)
    add_text(s, 84, top + 18, 400, 24, "可检索来源（候选集）", size=14, color=ACCENT, bold=True)
    sources = [
        ("window:latest_analysis", "最近一次成功观察"),
        ("window:summaries", "滚动结构化观察历史"),
        ("memory:working:observation", "当前工作观察卡片"),
        ("memory:items", "短期/稳定记忆项"),
        ("profile", "用户偏好与 persona"),
        ("assistant:chat:history", "当前运行期对话"),
        ("chat_history_fts", "跨会话 FTS5 索引"),
    ]
    for i, (k, v) in enumerate(sources):
        y = top + 52 + i * 28
        add_text(s, 96, y, 16, 20, "•", size=12, color=ACCENT)
        add_multi(s, 112, y, card_w - 60, 24, [
            ([(k, {"size": 12, "color": TEXT, "bold": True, "font": MONO}),
              ("　" + v, {"size": 12, "color": MUTED})], {}),
        ])

    # 右上：BM25 ranker
    add_round(s, SW - 60 - card_w, top, card_w, 220, fill=PANEL, line=LINE)
    add_text(s, SW - 60 - card_w + 24, top + 18, 400, 24, "ranker：SQLite FTS5 + BM25",
             size=14, color=ACCENT, bold=True)
    add_multi(s, SW - 60 - card_w + 24, top + 52, card_w - 48, 160, [
        ("内存 SQLite 对候选证据建临时 FTS5 索引", {"size": 12, "color": MUTED, "line_spacing": 1.5, "space_after": 4}),
        ([("中文 ", {"size": 12, "color": MUTED}),
          ("bigram 双字滑窗", {"size": 12, "color": TEXT, "bold": True}),
          ("分词，ASCII 整词保留", {"size": 12, "color": MUTED})],
         {"line_spacing": 1.5, "space_after": 4}),
        ("候选集通常 30-40 条，BM25 区分度充分", {"size": 12, "color": MUTED, "line_spacing": 1.5, "space_after": 4}),
        ([("零 LLM 调用", {"size": 12, "color": TEXT, "bold": True}),
          ("，毫秒级返回，确定性算法", {"size": 12, "color": MUTED})],
         {"line_spacing": 1.5}),
    ])

    # 右下：禁止的设计
    add_round(s, SW - 60 - card_w, top + 235, card_w, 145,
              fill=RGBColor(0x1A, 0x12, 0x14), line=RGBColor(0x5A, 0x2E, 0x2E))
    add_text(s, SW - 60 - card_w + 24, top + 253, 400, 24, "禁止的设计（已废弃）",
             size=14, color=BAD, bold=True)
    add_multi(s, SW - 60 - card_w + 24, top + 287, card_w - 48, 90, [
        ("每次固定返回 latest + recent + memory 大礼包", {"size": 12, "color": MUTED, "line_spacing": 1.45, "space_after": 3}),
        ("用 if \"当前\" in query 这类关键词判断", {"size": 12, "color": MUTED, "line_spacing": 1.45, "space_after": 3}),
        ("用 substring 命中判断情绪/意图/记忆相关性", {"size": 12, "color": MUTED, "line_spacing": 1.45, "space_after": 3}),
        ("用同一 VLM 做 reranker（连锁失败源）", {"size": 12, "color": MUTED, "line_spacing": 1.45}),
    ])


# =========================================================================
# 第 8 页 上下文分层
# =========================================================================
def slide_context():
    s = prs.slides.add_slide(BLANK)
    add_bg(s)
    add_header(s, "07 · 上下文管理", 8)
    add_title(s, "上下文分层 + prefix cache 冻结")
    add_lead(s, "不同来源的上下文进入方式不同，避免互相污染；profile 在会话级冻结以保证 llama.cpp prefix cache 命中率。")

    # 表格
    rows = [
        ("上下文层", "进入方式", "作用"),
        ("stable system prompt", "直接进入 messages", "约束回答原则与工具协议"),
        ("profile packet", "会话级冻结", "persona/profile 稳定，提升 prefix cache 命中"),
        ("session history", "最近 N 轮直接进入", "当前对话连续性"),
        ("window observations", "通过 memory.search 调取", "窗口/页面/代码/截图相关回答"),
        ("memory items", "近期尾部 + 检索候选", "用户偏好、任务事实"),
        ("cross-session chat", "chat_history_fts 检索", "找回历史讨论与结论"),
    ]
    top = 200
    row_h = 32
    col_x = [60, 360, 660]
    col_w = [300, 300, 560]
    # 表头
    for i, h in enumerate(rows[0]):
        add_rect(s, col_x[i], top, col_w[i], row_h, fill=RGBColor(0x14, 0x2A, 0x2E), line=LINE)
        add_text(s, col_x[i] + 14, top + 6, col_w[i] - 28, 22, h,
                 size=11, color=ACCENT, bold=True)
    # 数据行
    for r, row in enumerate(rows[1:], 1):
        y = top + r * row_h
        for i, cell in enumerate(row):
            add_rect(s, col_x[i], y, col_w[i], row_h, fill=PANEL, line=LINE)
            if i == 0:
                add_text(s, col_x[i] + 14, y + 6, col_w[i] - 28, 22, cell,
                         size=11, color=TEXT, bold=True, font=MONO)
            elif i == 1:
                add_text(s, col_x[i] + 14, y + 6, col_w[i] - 28, 22, cell,
                         size=11, color=MUTED)
            else:
                add_text(s, col_x[i] + 14, y + 6, col_w[i] - 28, 22, cell,
                         size=11, color=MUTED)

    # 底部三卡
    card_w, card_h = 380, 120
    top2 = 500
    cards = [
        ("会话级冻结",
         "ChatAgent 在首次 ask 时冻结 profile_packet 字符串，整个会话内复用同一对象。字节级一致是 KV cache 命中的前提。"),
        ("上下文预算边界",
         "预算只在模型调用边界生效，不在存储层生效。SQLite 不做破坏性字符截断——完整字段进库。"),
        ("跨会话检索",
         "对话结束时写入持久 chat_history_fts；下次会话只通过工具按需检索，不全量注入。"),
    ]
    for i, (t, body) in enumerate(cards):
        left = 60 + i * (card_w + 15)
        add_round(s, left, top2, card_w, card_h, fill=PANEL, line=LINE)
        add_text(s, left + 18, top2 + 14, card_w - 36, 22, t,
                 size=13, color=ACCENT, bold=True)
        add_text(s, left + 18, top2 + 44, card_w - 36, 70, body,
                 size=11, color=MUTED)


# =========================================================================
# 第 9 页 工程决策
# =========================================================================
def slide_decisions():
    s = prs.slides.add_slide(BLANK)
    add_bg(s)
    add_header(s, "08 · 工程决策", 9)
    add_title(s, "关键工程决策：删掉了什么、为什么")
    add_lead(s, "这个项目最有价值的不只是做了什么，更是明确删掉了什么。每一条删除都对应一次踩坑。")

    decisions = [
        ("废弃：独立 planner / orchestrator",
         "让一个模型先规划工具、另一个模型回答，链路长且不可控。改为单 ChatAgent 内部 probe → stream。"),
        ("废弃：多工具设计",
         "screen.look、memory.remember 等多工具让小 VLM 决策更不稳定。收敛为唯一工具 memory.search。"),
        ("废弃：LLM reranker",
         "用同一 MiniCPM-V 做 reranker 形成连锁失败。改用 BM25，毫秒级、确定性。"),
        ("废弃：关键词路由",
         "if \"左边\" in query 这类固定方位裁剪、substring 情绪判断——硬编码业务理解，脆弱且不可审计。改用 FTS5 BM25。"),
        ("废弃：默认注入窗口观察",
         "曾经把 latest_analysis 默认塞进对话 prompt。改为默认不注入，模型按需调用工具取证，避免上下文污染。"),
        ("废弃：增大 max_tokens 掩盖问题",
         "曾把 analyze_max_tokens 从 8192 提到 16384 \"解决\"截断。实际是 prompt 字段下限过高导致模型退化填充。回归 8192。"),
    ]
    top = 200
    card_w = (SW - 120 - 20) // 2
    card_h = 145
    for i, (t, body) in enumerate(decisions):
        col = i % 2
        row = i // 2
        left = 60 + col * (card_w + 20)
        y = top + row * (card_h + 12)
        add_round(s, left, y, card_w, card_h, fill=PANEL, line=LINE)
        add_text(s, left + 20, y + 14, card_w - 40, 26, t,
                 size=13, color=BAD, bold=True)
        add_text(s, left + 20, y + 48, card_w - 40, card_h - 60, body,
                 size=11, color=MUTED)


# =========================================================================
# 第 10 页 技术栈
# =========================================================================
def slide_stack():
    s = prs.slides.add_slide(BLANK)
    add_bg(s)
    add_header(s, "09 · 技术栈", 10)
    add_title(s, "技术栈：纯本地、零外部依赖")
    add_lead(s, "所有组件运行在本机，不依赖任何云服务或外部数据库。这是 local-first 隐私设计的硬约束。")

    rows = [
        ("层", "技术", "选型理由"),
        ("后端", "Python 3.11+ / FastAPI / Pydantic / uvicorn", "轻量、SSE 友好、与 ML 生态一致"),
        ("本地存储", "SQLite RuntimeStore / FTS5", "零外部 DB 依赖，单文件可备份"),
        ("检索", "SQLite FTS5 + BM25 / 中文 bigram", "毫秒级、确定性、无外部分词器依赖"),
        ("模型运行时", "llama.cpp llama-server / OpenAI API", "CPU/GPU 本地推理，协议兼容 OpenAI"),
        ("视觉模型", "MiniCPM-V 4.6 GGUF / 256K context", "开源小 VLM，支持截图输入"),
        ("桌面 UI", "Python + Win32 / 透明置顶悬浮窗", "原生低打扰，5 种状态机"),
        ("WebUI", "原生 HTML/CSS/JS / SSE 流式对话", "调试面板，可观测、可清空"),
        ("测试", "pytest / service-level + API route tests", "覆盖 ChatAgent / memory / window / API"),
    ]
    top = 200
    row_h = 32
    col_x = [60, 380, 820]
    col_w = [320, 440, 400]
    # 表头
    for i, h in enumerate(rows[0]):
        add_rect(s, col_x[i], top, col_w[i], row_h, fill=RGBColor(0x14, 0x2A, 0x2E), line=LINE)
        add_text(s, col_x[i] + 14, top + 6, col_w[i] - 28, 22, h,
                 size=11, color=ACCENT, bold=True)
    for r, row in enumerate(rows[1:], 1):
        y = top + r * row_h
        for i, cell in enumerate(row):
            add_rect(s, col_x[i], y, col_w[i], row_h, fill=PANEL, line=LINE)
            color = TEXT if i == 0 else MUTED
            bold = i == 0
            add_text(s, col_x[i] + 14, y + 6, col_w[i] - 28, 22, cell,
                     size=10, color=color, bold=bold)

    # pills
    pills = ["Windows 10+", "Python 3.11+", "FastAPI", "SQLite + FTS5",
             "MiniCPM-V 4.6", "llama.cpp", "Win32 API", "pytest"]
    top2 = 540
    x = 60
    for p in pills:
        w = len(p) * 9 + 28
        add_round(s, x, top2, w, 30, fill=RGBColor(0x14, 0x1F, 0x2E),
                  line=RGBColor(0x2E, 0x4A, 0x52), radius=0.5)
        add_text(s, x, top2 + 6, w, 22, p, size=11, color=ACCENT2,
                 align=PP_ALIGN.CENTER)
        x += w + 10


# =========================================================================
# 第 11 页 路线图
# =========================================================================
def slide_roadmap():
    s = prs.slides.add_slide(BLANK)
    add_bg(s)
    add_header(s, "10 · 进展与路线", 11)
    add_title(s, "当前进展与下一步")
    add_lead(s, "核心主线已全部落地。当前重点是把本地链路打磨稳定，而不是堆新功能。")

    items = [
        (True, "桌宠悬浮窗与 5 种状态（Idle / Observing / Analyzing / Privacy / Error）"),
        (True, "FastAPI 后端与 SQLite RuntimeStore"),
        (True, "MiniCPM-V 窗口结构化观察（含失败分型与 schema 收敛）"),
        (True, "memory.search(query) 单工具主线 + FTS5 BM25 ranker"),
        (True, "probe → stream 职责分离与流式工具调用循环"),
        (True, "跨会话对话 FTS5 索引 + runtime logs / tool traces / reset-all 调试接口"),
        (False, "工具结果按 token 预算裁剪"),
        (False, "观察质量评估（20 个真实窗口样本）+ 演示 GIF / 视频录制"),
        (False, "多显示器与窗口选择体验增强 · Rive 动画迁移"),
    ]
    top = 200
    row_h = 42
    for i, (done, t) in enumerate(items):
        y = top + i * row_h
        add_round(s, 60, y, SW - 120, row_h - 6, fill=PANEL, line=LINE)
        # 状态点
        dot_color = GOOD if done else WARN
        dot = s.shapes.add_shape(MSO_SHAPE.OVAL, PX(80), PX(y + 14), PX(14), PX(14))
        dot.fill.solid()
        dot.fill.fore_color.rgb = dot_color
        dot.line.fill.background()
        dot.shadow.inherit = False
        # 文字
        add_text(s, 110, y + 9, SW - 220, 24, t,
                 size=13, color=TEXT if not done else MUTED)
        # 状态标签
        label = "已完成" if done else "下一步"
        label_color = GOOD if done else WARN
        add_text(s, SW - 200, y + 9, 130, 24, label,
                 size=10, color=label_color, bold=True, align=PP_ALIGN.RIGHT)


# =========================================================================
# 第 12 页 运行方式
# =========================================================================
def slide_run():
    s = prs.slides.add_slide(BLANK)
    add_bg(s)
    add_header(s, "11 · 运行方式", 12)
    add_title(s, "如何跑起来：一键启动 + WebUI 调试")

    card_w, card_h = 600, 200
    top = 150
    # 左：启动
    add_round(s, 60, top, card_w, card_h, fill=PANEL, line=LINE)
    add_text(s, 84, top + 18, 300, 24, "一键启动", size=14, color=ACCENT, bold=True)
    add_code_block(s, 84, top + 52, card_w - 48, 100,
                   "# Windows PowerShell\ncd D:\\AI_Workspace\\window\n.\\scripts\\start_dev.cmd\n\n"
                   "# 启动前自检\npython .\\scripts\\check_environment.py --for-start",
                   size=10)
    add_text(s, 84, top + 165, card_w - 48, 22,
             "脚本会拉起 FastAPI 后端 + llama-server 模型服务 + 桌宠悬浮窗。",
             size=10, color=MUTED)

    # 右：WebUI
    add_round(s, SW - 60 - card_w, top, card_w, card_h, fill=PANEL, line=LINE)
    add_text(s, SW - 60 - card_w + 24, top + 18, 300, 24, "WebUI 调试入口",
             size=14, color=ACCENT, bold=True)
    add_code_block(s, SW - 60 - card_w + 24, top + 52, card_w - 48, 130,
                   "# 调试控制台\nhttp://127.0.0.1:18081/webui/\nhttp://127.0.0.1:18081/docs\n\n"
                   "# 关键调试 API\nGET  /api/webui/observations/latest\n"
                   "GET  /api/webui/tool-traces\nPOST /api/webui/reset-all",
                   size=10)

    # 底部典型流程
    top2 = 380
    add_round(s, 60, top2, SW - 120, 230, fill=RGBColor(0x12, 0x1F, 0x24),
              line=RGBColor(0x2E, 0x4A, 0x52))
    add_text(s, 84, top2 + 18, 300, 24, "典型运行流程",
             size=14, color=ACCENT, bold=True)

    steps = [
        ("1 · 启动", "start_dev.cmd → 后端 + 模型 + 浮窗"),
        ("2 · 自动观察", "watcher 后台截图 → VLM → SQLite"),
        ("3 · 用户提问", "点\"点击提问\" → 暂停观察"),
        ("4 · 取证回答", "probe → memory.search → stream"),
        ("5 · 调试复盘", "WebUI 看 trace / 截图 / logs"),
    ]
    n = len(steps)
    arrow_w = 18
    step_w = (SW - 168 - arrow_w * (n - 1)) // n
    step_h = 110
    s_top = top2 + 70
    for i, (t, body) in enumerate(steps):
        left = 84 + i * (step_w + arrow_w)
        fill = RGBColor(0x12, 0x1F, 0x24) if i == 0 else PANEL
        line = RGBColor(0x2E, 0x4A, 0x52) if i == 0 else LINE
        add_round(s, left, s_top, step_w, step_h, fill=fill, line=line)
        add_text(s, left + 12, s_top + 14, step_w - 24, 22, t,
                 size=12, color=TEXT, bold=True)
        add_text(s, left + 12, s_top + 42, step_w - 24, 60, body,
                 size=10, color=MUTED)
        if i < n - 1:
            add_text(s, left + step_w, s_top + 40, arrow_w, 30, "→",
                     size=20, color=ACCENT, align=PP_ALIGN.CENTER)


# =========================================================================
# 第 13 页 总结
# =========================================================================
def slide_end():
    s = prs.slides.add_slide(BLANK)
    add_bg(s)
    add_header(s, "12 · 总结", 13)

    # 居中标题
    add_multi(s, 140, 240, SW - 280, 130, [
        ([("三条主线，", {"size": 40, "color": TEXT, "bold": True}),
          ("一个工具", {"size": 40, "color": ACCENT, "bold": True})],
         {"align": PP_ALIGN.CENTER, "line_spacing": 1.2}),
    ])
    add_multi(s, 140, 340, SW - 280, 100, [
        ([("观察线产生证据 · 对话线按需取证 · 调试线全程可审计。", {"size": 14, "color": MUTED})],
         {"align": PP_ALIGN.CENTER, "line_spacing": 1.6}),
        ([("模型只看到一个工具 ", {"size": 14, "color": MUTED}),
          ("memory.search(query)", {"size": 14, "color": ACCENT2, "font": MONO, "bold": True}),
          ("，主链路不行就明确失败，不堆 fallback。", {"size": 14, "color": MUTED})],
         {"align": PP_ALIGN.CENTER, "line_spacing": 1.6}),
    ])
    add_multi(s, 140, 480, SW - 280, 100, [
        ("Local · Windows · Privacy-first", {"size": 12, "color": MUTED, "align": PP_ALIGN.CENTER, "space_after": 6}),
        ("Made by 宋林蔚 · github.com/fjnuslw", {"size": 12, "color": MUTED, "align": PP_ALIGN.CENTER, "space_after": 12}),
        ("谢谢，欢迎提问 →", {"size": 11, "color": MUTED, "align": PP_ALIGN.CENTER}),
    ])


# ---------- 生成 ----------
slide_cover()
slide_position()
slide_three_lines()
slide_arch()
slide_observe()
slide_dialog()
slide_tool()
slide_context()
slide_decisions()
slide_stack()
slide_roadmap()
slide_run()
slide_end()

out = "docs/Local_Window_Copilot.pptx"
prs.save(out)
print(f"saved: {out}")
print(f"slides: {len(prs.slides)}")
